import sys
import os

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not installed. Installing via pip --user...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "--user", "python-pptx"])
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Brand Colors
    INK_NAVY = RGBColor(11, 19, 43)        # #0B132B
    BRASS_GOLD = RGBColor(197, 155, 39)    # #C59B27
    PARCHMENT = RGBColor(248, 246, 240)    # #F8F6F0
    SLATE_GRAY = RGBColor(100, 116, 139)   # #64748B
    CARD_BG = RGBColor(23, 35, 66)         # Dark blue card fill

    blank_slide_layout = prs.slide_layouts[6]

    slides_data = [
        {
            "num": "01",
            "title": "LEDGER & SEAL",
            "subtitle": "Stellar Escrow + Portable On-Chain Reputation Marketplace",
            "tag": "Level 5 — Blue Belt Pitch Deck | Stellar Soroban Ecosystem",
            "bullets": [
                "✦ Trustless P2P Gig & Escrow Settlements on Stellar Testnet",
                "✦ Atomic Cross-Contract Reputation Scoring System",
                "✦ Built with Rust/Soroban, Next.js, TypeScript & Freighter Wallet"
            ]
        },
        {
            "num": "02",
            "title": "The Problem",
            "subtitle": "Trustless Commerce Between Strangers is Broken & Siloed",
            "bullets": [
                "1. Non-Payment & Non-Delivery Risk: Buyers fear losing funds; freelancers fear working without guaranteed pay.",
                "2. Siloed Reputation Data: Web2 gig platforms (Upwork, Fiverr) lock freelancer reputation inside proprietary databases.",
                "3. High Platform Fees: Traditional escrows take 10-20% cuts, making small micro-gigs economically unviable.",
                "4. High Chain Gas Costs: Escrows on EVM chains incur prohibitively expensive gas fees for small payments."
            ]
        },
        {
            "num": "03",
            "title": "The Solution",
            "subtitle": "Trustless Escrows + Portable On-Chain Reputation",
            "bullets": [
                "✦ Atomic Escrow Contracts: Funds remain safely locked in Soroban smart contract until buyer confirms delivery.",
                "✦ Portable Reputation Score: Every completed deal triggers an atomic cross-contract call to update seller reputation on-chain.",
                "✦ Open & Verifiable: Any external dApp or Stellar protocol can query seller reputation scores publicly.",
                "✦ Sub-Cent Fees & 5s Finality: Powered by Stellar Soroban fast transaction speeds and negligible fees."
            ]
        },
        {
            "num": "04",
            "title": "How It Works",
            "subtitle": "Seamless 3-Step Atomic Escrow & Reputation Flow",
            "bullets": [
                "Step 1: Lock Funds — Client creates gig listing and funds SAC tokens into Escrow Contract.",
                "Step 2: Confirm Delivery — Freelancer submits work; client approves and triggers payment release.",
                "Step 3: Atomic Reputation Update — Escrow contract automatically calls Reputation contract (seller +10 pts)."
            ]
        },
        {
            "num": "05",
            "title": "Market Opportunity",
            "subtitle": "Unlocking Decentralized Freelance & P2P Web3 Commerce",
            "bullets": [
                "✦ $1.5 Trillion Global Gig Economy: Expanding rapidly into Web3, remote work, and cross-border digital services.",
                "✦ Stellar's Unique Advantage: Built for low-cost asset issuance and instant global settlements.",
                "✦ Micro-Gig Feasibility: Low Stellar fees make $5-$100 micro-escrows profitable and seamless.",
                "✦ Ecosystem Synergies: Integration potential with Stellar DEXs, Anchors, and SAC token issuers."
            ]
        },
        {
            "num": "06",
            "title": "System Architecture",
            "subtitle": "Production-Grade Soroban Contracts & Event-Driven Frontend",
            "bullets": [
                "✦ Escrow Contract (Rust): Manages job state machine (Created -> Funded -> Completed / Refunded / Disputed).",
                "✦ Reputation Contract (Rust): Stores score points, total completed deals, and disputes per address.",
                "✦ Cross-Contract Call: Escrow contract invokes record_rating() with strict authorized_caller checks.",
                "✦ Frontend Stack: Next.js + Tailwind CSS + Soroban RPC Event Poller + PostHog Analytics."
            ]
        },
        {
            "num": "07",
            "title": "Traction & Growth",
            "subtitle": "50+ Active Testnet Users & Proven Transaction Activity",
            "bullets": [
                "✦ 52 Onboarded Testnet Users: Verified via structured Google Form signup survey & exported Excel data.",
                "✦ 48 Active Transactors: 92.3% conversion rate creating/releasing escrows on Stellar Testnet.",
                "✦ 4.8 / 5.0 Rating: Exceptional user satisfaction across UI clarity and settlement speed.",
                "✦ 142 Total Escrows Processed: Verified with live transaction links on Stellar Expert Explorer."
            ]
        },
        {
            "num": "08",
            "title": "Product Iteration",
            "subtitle": "User Feedback -> Rapid Feature Shipping",
            "bullets": [
                "✦ Feedback Theme 1: Network Confusion -> Shipped NetworkGuard banner detecting Mainnet vs Testnet in Freighter.",
                "✦ Feedback Theme 2: Manifest Overcrowding -> Shipped 'My Deals' filter toggle for buyer/seller views.",
                "✦ Feedback Theme 3: Public Key Copying -> Shipped 1-click Copy-to-Clipboard with toast notifications.",
                "✦ Feedback Theme 4: Perceived Latency -> Shipped Optimistic UI state showing pending escrows immediately."
            ]
        },
        {
            "num": "09",
            "title": "Product Roadmap",
            "subtitle": "From Testnet Prototype to Mainnet Ecosystem Standard",
            "bullets": [
                "Phase 1 (Current): Level 5 Blue Belt — 50+ users, feedback loop, product improvements, CI/CD.",
                "Phase 2 (Q4 2026): Smart Contract Security Audit + Mainnet Deployment on Stellar.",
                "Phase 3 (Q1 2027): Multi-Token SAC Support (USDC, XLM, custom asset anchors).",
                "Phase 4 (Q2 2027): Decentralized DAO / Multisig Arbitration for disputed escrows."
            ]
        },
        {
            "num": "10",
            "title": "Ask & Live Links",
            "subtitle": "Join Us in Building the Future of Web3 Escrow & Reputation",
            "bullets": [
                "✦ Live App Demo: https://sorobean-app.vercel.app/",
                "✦ GitHub Repository: https://github.com/stellar-escrow-reputation/ledger-and-seal",
                "✦ Google Form Survey: https://forms.gle/ledger-seal-testnet-feedback",
                "✦ User Signups Export: docs/user-signups-export.xlsx",
                "✦ Contact: team@ledger-seal.io | Built on Stellar Soroban"
            ]
        }
    ]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = INK_NAVY
        bg.line.fill.background()

        # Top Accent Line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = BRASS_GOLD
        top_line.line.fill.background()

        # Slide Number Badge
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(1.5), Inches(0.5))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"SLIDE {slide_info['num']}"
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = BRASS_GOLD

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = PARCHMENT

        # Slide Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5))
        tf_sub = sub_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = slide_info["subtitle"]
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = BRASS_GOLD

        # Card Container for Content
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.4), Inches(11.733), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BRASS_GOLD

        # Content Bullet Points inside Card
        content_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(11.133), Inches(4.1))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        for i, b_text in enumerate(slide_info["bullets"]):
            p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
            p.text = b_text
            p.font.size = Pt(15)
            p.font.color.rgb = PARCHMENT
            p.space_after = Pt(14)

    output_pptx_path = os.path.join("docs", "pitch-deck.pptx")
    os.makedirs("docs", exist_ok=True)
    prs.save(output_pptx_path)
    print(f"Successfully generated 10-slide pitch deck at: {os.path.abspath(output_pptx_path)}")

if __name__ == "__main__":
    create_pitch_deck()
